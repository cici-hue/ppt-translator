import streamlit as st
from pptx import Presentation
import openai
import re
import tempfile
import os
from pptx.util import Pt
from dotenv import load_dotenv
import uuid
# 新增：语言检测库
from langdetect import detect, LangDetectException
from langdetect.lang_detect_exception import LangDetectException

# ====================== 1. API Key Security Configuration (Unchanged) ======================
if os.path.exists(".env"):
    load_dotenv()
DEEPSEEK_API_KEY = st.secrets.get("DEEPSEEK_API_KEY") or os.getenv("DEEPSEEK_API_KEY")
if not DEEPSEEK_API_KEY:
    st.error("❌ DeepSeek API Key not configured! Please check environment variables or Streamlit Secrets.")
    st.stop()

client = openai.OpenAI(
    api_key=DEEPSEEK_API_KEY,
    base_url="https://api.deepseek.com"
)

# ====================== 2. Multi-Language Config + Font Mapping (Updated for English UI) ======================
# Language Config: {Display Name (English): (DeepSeek Standard Code, Language English Name, langdetect Code)}
# 新增：langdetect 对应的语言编码，确保检测结果能匹配
LANGUAGE_CONFIG = {
    "Chinese": ("zh", "Chinese", "zh-cn"),
    "English": ("en", "English", "en"),
    "German": ("de", "German", "de"),
    "Thai": ("th", "Thai", "th"),
    "Turkish": ("tr", "Turkish", "tr"),
    "Bengali": ("bn", "Bengali", "bn"),
    "Vietnamese": ("vi", "Vietnamese", "vi"),
    "French": ("fr", "French", "fr"),
    "Spanish": ("es", "Spanish", "es"),
    "Russian": ("ru", "Russian", "ru"),
    "Japanese": ("ja", "Japanese", "ja"),
    "Korean": ("ko", "Korean", "ko")
}
# Target Language - Font Mapping | Core: System-native fonts to avoid garbled text
FONT_MAP = {
    "zh": "微软雅黑",       # Target: Chinese
    "en": "Calibri",        # Target: English
    "de": "Calibri",        # Target: German
    "tr": "Calibri",        # Target: Turkish
    "fr": "Calibri",        # Target: French
    "es": "Calibri",        # Target: Spanish
    "ru": "Calibri",        # Target: Russian
    "th": "TH Sarabun New", # Target: Thai (Native to Windows/macOS)
    "vi": "VN Times",       # Target: Vietnamese (Native to Windows, compatible with Times New Roman on macOS)
    "bn": "Siyam Rupali",   # Target: Bengali (Native to Windows/macOS)
    "ja": "MS Mincho",      # Target: Japanese (Native to Windows)
    "ko": "Malgun Gothic"   # Target: Korean (Native to Windows)
}
# Extract language display names (for Streamlit dropdown)
LANG_NAMES = list(LANGUAGE_CONFIG.keys())

# ====================== 3. Utility Functions (Added Language Detection) ======================
def adjust_text_overflow_mild(text_frame, min_font_size=10):
    """Mild text overflow adjustment (Unchanged)"""
    if not text_frame or not text_frame.text.strip():
        return
    text_frame.word_wrap = True
    src_sizes = [run.font.size for para in text_frame.paragraphs for run in para.runs if run.font.size is not None]
    if not src_sizes:
        return
    current_font = src_sizes[0]
    for _ in range(6):
        try:
            if text_frame.height >= text_frame.text_height:
                break
        except:
            break
        new_font = current_font - Pt(1)
        new_font = new_font if new_font >= Pt(min_font_size) else Pt(min_font_size)
        for para in text_frame.paragraphs:
            for run in para.runs:
                if run.font.size is not None:
                    run.font.size = new_font
        current_font = new_font
    if current_font == Pt(min_font_size):
        try:
            if text_frame.height < text_frame.text_height:
                st.warning(f"💡 Some text is slightly overflowing (reduced to minimum 10pt). Please manually adjust text box width.")
        except:
            pass

def detect_text_language(text):
    """新增：检测文本语言，返回langdetect编码，失败则返回None"""
    if not text or len(text.strip()) < 2:  # 文本过短无法准确检测
        return None
    try:
        return detect(text.strip())
    except LangDetectException:
        return None

def translate_text(text, src_lang_code, src_lang_name, tgt_lang_code, tgt_lang_name):
    """Multi-language translation function | Pass language code + name (English UI adaptation)"""
    if not text or not text.strip():
        return text
    try:
        # Dynamically adapt translation prompt to source/target language
        system_prompt = f"""You are a professional multilingual translation expert, proficient in {src_lang_name} and {tgt_lang_name} mutual translation. Strictly follow these rules:
1. Terminology accuracy: Use industry-standard translations for business/office PPT professional terms, maintain consistency;
2. Format preservation: Keep line breaks, spaces, punctuation, numbers/units in the original text unchanged, no additions or deletions;
3. Expression adaptation: Conform to PPT reading habits of the target language, concise and powerful titles, fluent text;
4. No extra output: Only return translation results, no explanations, notes, punctuation corrections or irrelevant content;
5. Special characters: Accurately handle special characters/diacritics of the target language (e.g., German umlauts, Vietnamese tones)."""
        response = client.chat.completions.create(
            model="deepseek-chat",
            messages=[{"role": "system", "content": system_prompt}, {"role": "user", "content": text}],
            temperature=0.1,  # Low temperature ensures stable translation results
            max_tokens=3000    # Increase token limit for multi-language long text
        )
        return response.choices[0].message.content.strip()
    except Exception as e:
        st.error(f"❌ Translation error: {str(e)}")
        return text

def translate_ppt(input_file_path, output_file_path, src_lang, tgt_lang):
    """Core PPT translation logic | Added language detection to skip target language text"""
    # Parse source/target language config (code + name + langdetect code)
    src_lang_code, src_lang_name, src_detect_code = LANGUAGE_CONFIG[src_lang]
    tgt_lang_code, tgt_lang_name, tgt_detect_code = LANGUAGE_CONFIG[tgt_lang]
    # Dynamically match target font (solve garbled text)
    target_font = FONT_MAP[tgt_lang_code]
    
    try:
        prs = Presentation(input_file_path)
        st.success(f"✅ PPT loaded successfully | Total slides: {len(prs.slides)} | Source Language: {src_lang} | Target Language: {tgt_lang} | Adapted Font: {target_font}")
    except Exception as e:
        st.error(f"❌ Failed to load PPT: {str(e)}")
        return False
    
    total_texts, translated_texts, skipped_texts = 0, 0, 0  # 新增：skipped_texts 统计跳过的文本数
    # Progress bar + status prompt (English adaptation)
    progress_bar = st.progress(0)
    status_text = st.empty()

    for slide_idx, slide in enumerate(prs.slides, 1):
        status_text.text(f"🔄 Processing slide {slide_idx}/{len(prs.slides)}...")
        progress_bar.progress(slide_idx / len(prs.slides))

        for shape in slide.shapes:
            # Process text boxes (multi-language font adaptation, format preservation)
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    original_text = paragraph.text.strip()
                    if original_text:
                        total_texts += 1
                        # 新增：检测文本语言，若已是目标语言则跳过
                        detected_lang = detect_text_language(original_text)
                        if detected_lang == tgt_detect_code:
                            skipped_texts += 1
                            continue  # 跳过目标语言文本，不翻译
                        
                        # Call multi-language translation function
                        translated_text = translate_text(original_text, src_lang_code, src_lang_name, tgt_lang_code, tgt_lang_name)
                        if translated_text and translated_text != original_text:
                            # Preserve original format (bold/color/font size)
                            src_font = paragraph.runs[0].font if paragraph.runs else None
                            paragraph.text = ""
                            new_run = paragraph.add_run()
                            new_run.text = translated_text
                            if src_font:
                                new_run.font.bold = src_font.bold if src_font.bold is not None else False
                                new_run.font.size = src_font.size
                                new_run.font.name = target_font  # Multi-language font adaptation
                                try:
                                    new_run.font.color.rgb = src_font.color.rgb
                                except:
                                    pass
                            # 1x line spacing (supported in all versions)
                            paragraph.line_spacing = 1
                            # Mild overflow adjustment
                            adjust_text_overflow_mild(shape.text_frame)
                            translated_texts += 1
            # Process tables (same as text boxes, multi-language font + format preservation)
            if shape.has_table:
                try:
                    table = shape.table
                    for row_idx, row in enumerate(table.rows):
                        for cell_idx, cell in enumerate(row.cells):
                            cell_text = cell.text.strip()
                            if cell_text:
                                total_texts += 1
                                # 新增：检测表格文本语言，若已是目标语言则跳过
                                detected_lang = detect_text_language(cell_text)
                                if detected_lang == tgt_detect_code:
                                    skipped_texts += 1
                                    continue
                                
                                translated_cell = translate_text(cell_text, src_lang_code, src_lang_name, tgt_lang_code, tgt_lang_name)
                                if translated_cell and translated_cell != cell_text:
                                    cell_src_font = None
                                    if cell.text_frame.paragraphs and cell.text_frame.paragraphs[0].runs:
                                        cell_src_font = cell.text_frame.paragraphs[0].runs[0].font
                                    cell.text_frame.clear()
                                    cell_para = cell.text_frame.add_paragraph()
                                    cell_run = cell_para.add_run()
                                    cell_run.text = translated_cell
                                    # Format preservation + multi-language font
                                    if cell_src_font:
                                        cell_run.font.bold = cell_src_font.bold if cell_src_font.bold is not None else False
                                        cell_run.font.size = cell_src_font.size
                                        cell_run.font.name = target_font
                                        try:
                                            cell_run.font.color.rgb = cell_src_font.color.rgb
                                        except:
                                            pass
                                    cell_para.line_spacing = 1
                                    adjust_text_overflow_mild(cell.text_frame)
                                    translated_texts += 1
                except Exception as e:
                    st.warning(f"⚠️ Table processing exception (skipped): {str(e)[:40]}...")

    # Save translated PPT (English adaptation)
    try:
        prs.save(output_file_path)
        progress_bar.progress(100)
        status_text.text("✅ Translation completed!")
        # Multi-language translation statistics (Updated: add skipped count)
        st.success(f"""
        📊 Translation Statistics | Source Language: {src_lang} → Target Language: {tgt_lang}
        ├─ Total text blocks (text boxes + tables): {total_texts}
        ├─ Successfully translated text blocks: {translated_texts}
        ├─ Skipped target language text blocks: {skipped_texts}  # 新增：显示跳过的文本数
        ├─ Target language adapted font: {target_font}
        └─ Format preservation: 1:1 retention of bold/color/font size + 1x line spacing + mild overflow adjustment
        """)
        return True
    except Exception as e:
        st.error(f"❌ Failed to save PPT: {str(e)} (Please close the target PPT file and try again)")
        return False

# ====================== 4. Streamlit Web Interface (Fully English) ======================
def main():
    st.set_page_config(page_title="PPT Smart Translation Tool", page_icon="📄", layout="wide")
    st.title("📄 PPT Smart Translation Tool")
    st.divider()

    # Sidebar: Translation Configuration (Fully English)
    with st.sidebar:
        st.header("⚙️ Translation Settings")
        # Multi-language source language selection (default: Chinese)
        src_lang = st.selectbox("🔤 Source Language", LANG_NAMES, index=LANG_NAMES.index("Chinese"))
        # Multi-language target language selection (default: English)
        tgt_lang = st.selectbox("🌐 Target Language", LANG_NAMES, index=LANG_NAMES.index("English"))
        # Validation: Source language != Target language
        if src_lang == tgt_lang:
            st.error("❌ Source language and target language cannot be the same! Please reselect.")
            st.stop()
        # Feature description (Updated: add language detection feature)
        st.info("""
        📌 Core Features
        1. Supports mutual translation between 12 mainstream languages;
        2. Automatically adapts target language fonts to avoid garbled text;
        3. Preserves all original PPT formats;
        4. Supports text box/table translation;
        5. Only supports .pptx format - upload file, translate with one click, download results;
        6. Smart language detection: Skip translation for text already in target language (new feature).
        """)
        st.warning("""
        ⚠️ Important Notes
        1. Recommended to upload PPT files smaller than 20MB for faster translation;
        2. Complex artistic text/special shape text may not be parsed (python-pptx library limitation);
        3. Please verify professional terminology in translation results to ensure accuracy;
        4. Language detection requires at least 2 characters (short text may not be detected accurately).
        """)

    # Main Interface: File Upload (Fully English)
    st.subheader("📤 Upload PPT File (Only .pptx format supported)")
    uploaded_file = st.file_uploader("Click to select or drag PPT file here", type=["pptx"], accept_multiple_files=False)

    if uploaded_file is not None:
        # Display uploaded file information (English)
        file_size = round(uploaded_file.size / 1024 / 1024, 2)
        st.info(f"📁 Uploaded File: {uploaded_file.name} | File Size: {file_size} MB")
        # Generate unique temporary file name (avoid conflicts)
        unique_id = str(uuid.uuid4())[:8]
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as temp_input:
            temp_input.write(uploaded_file.getbuffer())
            temp_input_path = temp_input.name

        # Translation button (prominent, English)
        if st.button("🚀 Start Multi-Language Translation", type="primary", use_container_width=True):
            # Generate output temporary file
            temp_output_path = os.path.join(tempfile.gettempdir(), f"ppt_translated_{unique_id}.pptx")
            # Execute multi-language translation
            translate_success = translate_ppt(temp_input_path, temp_output_path, src_lang, tgt_lang)
            # Provide download link (dynamically generated file name, e.g., "filename_Chinese_to_German.pptx")
            if translate_success and os.path.exists(temp_output_path):
                download_file_name = f"{os.path.splitext(uploaded_file.name)[0]}_{src_lang}_to_{tgt_lang}.pptx"
                with open(temp_output_path, "rb") as f:
                    st.download_button(
                        label="📥 Download Translated PPT File",
                        data=f,
                        file_name=download_file_name,
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        type="primary",
                        use_container_width=True
                    )
            # Clean up temporary files (avoid disk usage)
            os.unlink(temp_input_path)
            if os.path.exists(temp_output_path):
                os.unlink(temp_output_path)

if __name__ == "__main__":
    main()
